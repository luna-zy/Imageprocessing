from pptx import Presentation
from pptx.util import Inches

# 创建PPT
prs = Presentation()

# 封面页
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "BRISK: Short and Long Distance Pairs"
subtitle.text = "Understanding BRISK Feature Descriptor\nby Zhang Zhenhao"

# 目录页
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Outline"
content = slide.placeholders[1]
content.text = "1. Introduction to BRISK\n2. Short and Long Distance Pairs\n3. Why Use Both?\n4. Mathematical Formulation\n5. Implementation in OpenCV\n6. Summary"

# BRISK 简介
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Introduction to BRISK"
content = slide.placeholders[1]
content.text = "• BRISK (Binary Robust Invariant Scalable Keypoints) is a feature detection and description algorithm.\n"
content.text += "• It extracts keypoints from an image and represents them as binary descriptors.\n"
content.text += "• Used in real-time applications like SLAM and image matching."

# 短距离像素对
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Short Distance Pairs"
content = slide.placeholders[1]
content.text = "• Short distance pairs are adjacent pixels around a keypoint.\n"
content.text += "• Used to compute the keypoint’s dominant orientation.\n"
content.text += "• This ensures rotation invariance.\n\n"
content.text += "Formula:\nθ = tan⁻¹ ( Σ w_i (y_i - y_c) / Σ w_i (x_i - x_c) )"

# 长距离像素对
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Long Distance Pairs"
content = slide.placeholders[1]
content.text = "• Long distance pairs are pixels farther apart on the sampling rings.\n"
content.text += "• Used to construct the binary descriptor.\n"
content.text += "• Comparison of pixel intensities generates a 512-bit descriptor.\n\n"
content.text += "Descriptor calculation:\n If I(p1) < I(p2), set bit = 1, else set bit = 0."

# 为什么使用两种像素对
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Why Use Both?"
content = slide.placeholders[1]
content.text = "• Short Distance Pairs → Compute Orientation (Rotation Invariance)\n"
content.text += "• Long Distance Pairs → Generate Binary Descriptor (Feature Matching)\n"
content.text += "• Combination makes BRISK robust in feature detection and matching."

# 数学公式
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Mathematical Formulation"
content = slide.placeholders[1]
content.text = "• Orientation Calculation:\nθ = tan⁻¹ ( Σ w_i (y_i - y_c) / Σ w_i (x_i - x_c) )\n\n"
content.text += "• Descriptor Calculation:\nBinary vector based on pixel intensity comparison."

# OpenCV 实现
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Implementation in OpenCV"
content = slide.placeholders[1]
content.text = "```python\nimport cv2\n\nimage = cv2.imread('image.jpg', cv2.IMREAD_GRAYSCALE)\nbrisk = cv2.BRISK_create()\nkeypoints, descriptors = brisk.detectAndCompute(image, None)\n```"

# 总结
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Summary"
content = slide.placeholders[1]
content.text = "• BRISK uses a combination of Short and Long Distance Pairs.\n"
content.text += "• Short Distance Pairs help in rotation invariance.\n"
content.text += "• Long Distance Pairs build the binary descriptor.\n"
content.text += "• This makes BRISK fast, efficient, and robust for feature matching."

# 保存PPT
pptx_filename = "/mnt/data/BRISK_Short_Long_Pairs.pptx"
prs.save(pptx_filename)

# 提供下载链接
pptx_filename
